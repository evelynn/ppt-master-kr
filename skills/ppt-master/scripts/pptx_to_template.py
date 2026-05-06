#!/usr/bin/env python3
"""
PPT Master - PPTX → DESIGN.md + layout template (v1, framework only).

Reads a sample .pptx and emits a starter
templates/layouts/<slug>/{DESIGN.md, design_spec.md, source_meta.json}.
SVG slide stubs are populated as freehex slides; only theme tokens are
extracted automatically. Per-vendor extractors live in `_extractors/<slug>.py`
and run after the generic pass to override colors / layout / components.

Per the project plan, full PPTX import is staged: v1 covers a generic PPTX
(theme colors + master fonts + a list of slide layouts) so that the rest of
the toolchain can be exercised. v2 work happens once a real DRM-cleared
sample arrives — see Phase E in the plan.

Usage:
    python3 pptx_to_template.py <sample.pptx> --name <slug>
        [--style brand|general|consultant]
        [--fallback raster]            # raster the layouts as PNGs into preview/

Output:
    templates/layouts/<slug>/
        DESIGN.md
        design_spec.md
        source_meta.json
        preview/*.png        (only with --fallback raster)
"""

from __future__ import annotations

import argparse
import datetime
import hashlib
import importlib
import json
import sys
from collections import Counter
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError as e:
    Presentation = None  # type: ignore
    _PPTX_IMPORT_ERROR = str(e)


SCRIPTS = Path(__file__).parent
SKILL_ROOT = SCRIPTS.parent
LAYOUTS_DIR = SKILL_ROOT / "templates" / "layouts"


# ---------------------------------------------------------------------------
# Theme color and font extraction
# ---------------------------------------------------------------------------

def _hex_from_color(color) -> str | None:
    try:
        rgb = color.rgb
        if rgb is None:
            return None
        return "#" + str(rgb).upper()
    except Exception:
        return None


def extract_colors(prs) -> dict[str, str]:
    """Best-effort palette extraction from theme + slide masters."""
    palette: dict[str, str] = {}
    # Slide master color scheme (most reliable)
    for master in prs.slide_masters:
        try:
            scheme = master.element.findall(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}clrScheme"
            )
            if scheme:
                _ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                for child in scheme[0]:
                    role = child.tag[len(_ns):]  # dk1 / lt1 / accent1..6 / hlink / folHlink
                    srgb = child.find(f"{_ns}srgbClr")
                    if srgb is not None:
                        palette[role] = "#" + srgb.attrib["val"].upper()
        except Exception:
            continue

    # Sample frequent fills from body slides
    fills = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                fill = shape.fill
                if getattr(fill, "type", None) is not None and fill.type == 1:  # solid
                    h = _hex_from_color(fill.fore_color)
                    if h:
                        fills[h] += 1
            except Exception:
                continue
    # Top 4 fills get brand-1..brand-4 (skip already-known ones)
    known_hex = {v.upper() for v in palette.values()}
    used = []
    for hexv, _count in fills.most_common(20):
        if hexv.upper() in known_hex:
            continue
        used.append(hexv)
        if len(used) >= 4:
            break

    # Map clrScheme roles to our token names.
    rename = {
        "dk1": "primary",
        "lt1": "canvas",
        "dk2": "ink",
        "lt2": "surface-1",
        "accent1": "brand-1",
        "accent2": "brand-2",
        "accent3": "accent",
        "accent4": "brand-3",
        "accent5": "brand-4",
        "accent6": "brand-5",
    }
    out: dict[str, str] = {}
    for role, hexv in palette.items():
        out[rename.get(role, role)] = hexv
    for i, hexv in enumerate(used):
        out.setdefault(f"brand-extra-{i+1}", hexv)
    out.setdefault("ink-muted", "#6B6B73")
    out.setdefault("surface-2", "#E5E5EA")
    out.setdefault("ink-inverse", "#FFFFFF")
    return out


def extract_fonts(prs) -> dict[str, str]:
    fonts = {"heading": "DM Sans, Pretendard, sans-serif",
             "body": "Inter, Pretendard, sans-serif",
             "code": "JetBrains Mono, Consolas, monospace"}
    for master in prs.slide_masters:
        try:
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            major = master.element.find(f".//{ns}majorFont/{ns}latin")
            minor = master.element.find(f".//{ns}minorFont/{ns}latin")
            if major is not None:
                fonts["heading"] = f'"{major.attrib["typeface"]}", Pretendard, sans-serif'
            if minor is not None:
                fonts["body"] = f'"{minor.attrib["typeface"]}", Pretendard, sans-serif'
            break
        except Exception:
            continue
    return fonts


def extract_typography_distribution(prs) -> list[float]:
    sizes = Counter()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        sizes[round(run.font.size.pt)] += 1
    return [s for s, _ in sizes.most_common(8)]


# ---------------------------------------------------------------------------
# DESIGN.md emission
# ---------------------------------------------------------------------------

def render_design_md(*, slug: str, style: str, colors: dict[str, str],
                      fonts: dict[str, str], font_sizes_pt: list[float],
                      canvas_viewbox: str) -> str:
    color_rows: list[str] = []
    for tok, hexv in colors.items():
        color_rows.append(f"| {tok} | {hexv} | extracted | from PPTX |")

    type_steps = ["hero-display", "heading-xl", "heading-lg",
                   "heading-md", "body-lg", "body-md", "micro"]
    type_rows = []
    sizes_pt = list(reversed(sorted(font_sizes_pt))) if font_sizes_pt else [54, 36, 28, 22, 18, 14, 10]
    while len(sizes_pt) < len(type_steps):
        sizes_pt.append(max(8, sizes_pt[-1] - 4))
    for tok, sz_pt in zip(type_steps, sizes_pt[:len(type_steps)]):
        sz_px = round(float(sz_pt) * 1.333)
        weight = 700 if "heading" in tok or "hero" in tok else 400
        type_rows.append(
            f"| {tok} | {sz_px} | {weight} | 1.30 | 0 | extracted from PPTX |"
        )

    return f"""# {slug} — DESIGN.md

> Auto-extracted from a sample PPTX on {datetime.date.today().isoformat()}.
> Review and tighten section by section. Section headings are parsed; do not
> rename. See `references/design-md-spec.md` for the full specification.

## Overview

Imported sample style: {style}. Theme colors and master fonts were extracted
from the slide master; type sizes were inferred from text-run frequency.
Hand-tune any values flagged in **Known Gaps**.

## Colors

### Brand & Accent
| Token | Hex | Role | Usage |
| ----- | --- | ---- | ----- |
{chr(10).join(color_rows)}

### Surface
| Token | Hex | Role |
| ----- | --- | ---- |

### Text
| Token | Hex | Role |
| ----- | --- | ---- |

### Semantic
| Token | Hex | Role |
| ----- | --- | ---- |

## Typography

**Heading family**: {fonts.get("heading", "")}
**Body family**: {fonts.get("body", "")}
**Code family**: {fonts.get("code", "")}

### Hierarchy
| Token | Size (px) | Weight | Line height | Letter spacing | Use |
| ----- | --------- | ------ | ----------- | -------------- | --- |
{chr(10).join(type_rows)}

## Layout

### Canvas
| Format | viewBox | Margins (T/R/B/L) | Grid columns | Gutter |
| ------ | ------- | ----------------- | ------------ | ------ |
| ppt169 | {canvas_viewbox} | 60 / 80 / 60 / 80 | 12 | 24 |

### Spacing scale
| Token | px |
| ----- | -- |
| xxs | 4 |
| xs  | 8 |
| sm  | 12 |
| md  | 16 |
| lg  | 24 |
| xl  | 32 |
| xxl | 48 |
| 3xl | 64 |
| 4xl | 96 |

## Elevation & Depth

| Token | Effect | SVG implementation |
| ----- | ------ | ------------------ |
| elevation.0 | Flat | (no filter) |
| elevation.1 | Subtle | `filter="url(#shadow-1)"` |

## Shapes

### Rounded corner scale
| Token | px |
| ----- | -- |
| xs  | 4 |
| sm  | 8 |
| md  | 12 |
| lg  | 16 |
| xl  | 24 |
| hero | 32 |
| full | 9999 |

## Components

### product-cards/coral
- Auto-imported template uses palette token `{{colors.brand-1}}`. Hand-tune.

## Slide Templates

### cover_default
- Imported placeholder cover. Edit by hand.

### content_default
- Imported placeholder content layout.

## Do's and Don'ts

### Do
- Reference colors via `{{colors.<token>}}`, not raw hex.
- Verify the extracted palette before generating downstream slides.

### Don't
- Don't trust the extracted hierarchy without inspection.

## Canvas Variants

| Element | ppt169 |
| ------- | ------ |
| spacing.xl | 32 |

## Known Gaps

- Source-specific shape vocabulary, SmartArt, charts, and master placeholder
  geometry have not been imported. Add per-template extractor in
  `scripts/_extractors/{slug}.py`.
"""


def write_template(prs, *, slug: str, style: str, fallback: str | None,
                    output_root: Path) -> Path:
    target = output_root / slug
    target.mkdir(parents=True, exist_ok=True)
    (target / "preview").mkdir(exist_ok=True)

    colors = extract_colors(prs)
    fonts = extract_fonts(prs)
    sizes_pt = extract_typography_distribution(prs)

    sw = prs.slide_width
    sh = prs.slide_height
    # Convert EMU → px (assume 96 dpi: 1 px = 9525 EMU). Cap to 1280x720 for ppt169.
    px_w = round(sw / 9525)
    px_h = round(sh / 9525)
    canvas_viewbox = f"0 0 {px_w} {px_h}"

    design_md = render_design_md(
        slug=slug, style=style, colors=colors, fonts=fonts,
        font_sizes_pt=sizes_pt, canvas_viewbox=canvas_viewbox,
    )
    (target / "DESIGN.md").write_text(design_md, encoding="utf-8")

    # Stub slide SVGs: one rectangle filled with primary, hand-edit later.
    slide_stubs = {
        "01_cover.svg": '<rect x="0" y="0" width="100%" height="100%" fill="{colors.primary}"/>',
        "02_chapter.svg": '<rect x="0" y="0" width="100%" height="100%" fill="{colors.canvas}"/>',
        "03_content.svg": '<rect x="0" y="0" width="100%" height="100%" fill="{colors.canvas}"/>',
        "04_ending.svg": '<rect x="0" y="0" width="100%" height="100%" fill="{colors.primary}"/>',
    }
    for name, body in slide_stubs.items():
        (target / name).write_text(
            f'<svg xmlns="http://www.w3.org/2000/svg" viewBox="{canvas_viewbox}">\n  {body}\n</svg>\n',
            encoding="utf-8",
        )

    (target / "design_spec.md").write_text(
        f"# {slug} - design_spec.md\n\n"
        f"Auto-stub. Replace placeholder SVGs and tighten DESIGN.md before use.\n",
        encoding="utf-8",
    )

    # Per-template extractor hook
    extractor_module = f"_extractors.{slug}"
    try:
        sys.path.insert(0, str(SCRIPTS))
        mod = importlib.import_module(extractor_module)
        if hasattr(mod, "refine"):
            mod.refine(prs, target=target)
            print(f"  • per-template extractor {extractor_module} applied")
    except ImportError:
        pass

    if fallback == "raster":
        try:
            # Lazy: just record that fallback was requested; rasterization
            # belongs to libreoffice / unoconv and is out of scope for v1.
            (target / "preview" / "README.txt").write_text(
                "Add raster previews here (e.g. via libreoffice --convert-to png).\n",
                encoding="utf-8",
            )
        except Exception:
            pass

    src_path = Path(prs.element.attrib.get("name") or "<sample.pptx>")
    (target / "source_meta.json").write_text(
        json.dumps({
            "imported_at": datetime.datetime.now(datetime.timezone.utc).isoformat(),
            "extractor_version": 1,
            "style_hint": style,
            "extracted_palette_size": len(colors),
            "extracted_font_sizes_pt": sizes_pt,
        }, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )

    return target


# ---------------------------------------------------------------------------
# CLI
# ---------------------------------------------------------------------------

def main() -> int:
    p = argparse.ArgumentParser(description=__doc__.split("\n\n")[0])
    p.add_argument("pptx", type=Path)
    p.add_argument("--name", required=True, help="Template slug (e.g. acme_brand)")
    p.add_argument("--style", default="general", choices=["brand", "general", "consultant", "scenario", "government", "special"])
    p.add_argument("--fallback", choices=["raster"], default=None)
    p.add_argument("--output-root", type=Path, default=LAYOUTS_DIR)
    args = p.parse_args()

    if Presentation is None:
        print("✖ python-pptx is required: pip install python-pptx")
        return 1
    if not args.pptx.exists():
        print(f"✖ PPTX not found: {args.pptx}")
        return 1

    prs = Presentation(str(args.pptx))
    target = write_template(
        prs, slug=args.name, style=args.style,
        fallback=args.fallback, output_root=args.output_root,
    )

    # Update layouts_index.json (best-effort).
    index_path = args.output_root / "layouts_index.json"
    if index_path.exists():
        try:
            idx = json.loads(index_path.read_text(encoding="utf-8"))
            cat = idx.setdefault("categories", {}).setdefault(args.style, {"label": args.style, "layouts": []})
            if args.name not in cat.setdefault("layouts", []):
                cat["layouts"].append(args.name)
            idx.setdefault("layouts", {})[args.name] = {
                "label": args.name,
                "summary": f"Imported from {args.pptx.name}",
                "tone": "to-be-defined",
                "themeMode": "Auto-extracted",
                "keywords": ["imported"],
            }
            idx.setdefault("meta", {})["total"] = len(idx.get("layouts", {}))
            index_path.write_text(json.dumps(idx, indent=2, ensure_ascii=False), encoding="utf-8")
        except Exception as e:
            print(f"  ⚠ failed to update layouts_index.json: {e}")

    print(f"✓ wrote {target}")
    print(f"  next: edit DESIGN.md, replace stub SVGs, then `python3 design_tokens.py validate {target}/DESIGN.md`")
    return 0


if __name__ == "__main__":
    sys.exit(main())
